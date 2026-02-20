import os
import requests
import textwrap
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK

def fetch_image(query):
    """
    Fetches an image from Pexels based on the query.
    Returns the local path to the downloaded image or None if failed.
    """
    api_key = os.getenv("PEXELS_API_KEY")
    if not api_key:
        print("Warning: PEXELS_API_KEY not set.")
        return None

    try:
        # Force landscape orientation and large size to avoid random vertical/weird images
        url = f"https://api.pexels.com/v1/search?query={query}&per_page=15&orientation=landscape&size=medium"
        headers = {
            "Authorization": api_key
        }
        
        response = requests.get(url, headers=headers, timeout=10)
        
        if response.status_code == 200:
            data = response.json()
            if data['photos']:
                # Get medium size image
                image_url = data['photos'][0]['src']['medium']
                
                # Download image
                img_response = requests.get(image_url, timeout=10)
                if img_response.status_code == 200:
                    filename = f"temp_{query.replace(' ', '_')}.jpg"
                    # Sanitize filename
                    filename = "".join([c for c in filename if c.isalpha() or c.isdigit() or c in (' ', '.', '_')]).strip()
                    path = os.path.join("static", filename)
                    
                    # Ensure static dir exists
                    if not os.path.exists("static"):
                        os.makedirs("static")
                        
                    with open(path, 'wb') as f:
                        f.write(img_response.content)
                    return path
            else:
                print(f"No results found on Pexels for: {query}")
        else:
            print(f"Pexels API Error: {response.status_code} - {response.text}")
            
    except Exception as e:
        print(f"Error fetching image for '{query}': {e}")
    return None

from datetime import datetime

def fit_text(text_frame, text_list, font_size_pt=22, max_height_inches=5.0, is_quote=False, font_family='Arial'):
    """
    Adds text to text_frame and reduces font size until it fits.
    Applies professional styling.
    """
    current_size = font_size_pt
    min_size = 9 # Allow going smaller
    
    apply_tight_spacing = False
    
    # helper to check fit
    def check_fit(size, tight_spacing=False):
        # Precise Height Calculation
        # Width in points
        try:
            box_width_pt = text_frame._parent.width / 12700 
        except:
            box_width_pt = 500 # Fallback 
            
        # Heuristic for Arial char width (Conservative: 0.6)
        avg_char_width_pt = size * 0.6 
        chars_per_line = box_width_pt / avg_char_width_pt
        if chars_per_line < 1: chars_per_line = 1
        
        spacing_mult = 1.1 if tight_spacing else 1.3
        
        total_height_pt = 0
        for point in text_list:
            # Bullet + Text
            length = len(point) + 4
            lines = (length // chars_per_line) + 1
            # Height = lines * line_height + space_after
            para_h = (lines * size * spacing_mult) + (size * 0.5)
            total_height_pt += para_h
            
        return (total_height_pt / 72) <= (max_height_inches - 0.2) # 0.2 buffer

    while current_size >= min_size:
        # Check normal spacing
        if current_size > 14:
            if check_fit(current_size, tight_spacing=False):
                 apply_tight_spacing = False
                 break
        else:
            # For small fonts, check if we need tight spacing
            if check_fit(current_size, tight_spacing=False):
                apply_tight_spacing = False
                break
            # Try tight spacing
            if check_fit(current_size, tight_spacing=True):
                apply_tight_spacing = True
                break
                
        current_size -= 1 # Granular decrease
    
    if current_size < min_size:
        current_size = min_size
        apply_tight_spacing = True
        print("Warning: Text might still overflow.")

    # Apply final text
    text_frame.clear()
    text_frame.word_wrap = True
    
    for point in text_list:
        p = text_frame.add_paragraph()
        clean_point = point.strip()
        if not is_quote and clean_point and not clean_point.startswith("•") and not clean_point.startswith("-"):
            p.text = "• " + clean_point
        else:
            p.text = clean_point
            
        p.font.name = font_family
        p.font.size = Pt(current_size)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = Pt(current_size * 0.5)
        p.line_spacing = 1.1 if apply_tight_spacing else 1.3
        
        if is_quote:
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        else:
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT


def get_smart_font_size(text_length, layout_type):
    """
    Returns a starting font size based on text length and layout.
    """
    if layout_type == 'big_quote':
        return 36
        
    if text_length < 200:
        return 24
    elif text_length < 400:
        return 18
    elif text_length < 600:
        return 14
    else:
        return 12


def create_ppt(presentation_data, theme_color=(41, 128, 185), base_ppt=None, background_image=None):
    """
    Create PPT from JSON data dictionary.
    base_ppt: Path to an existing PPTX to append slides to (preserves theme).
    background_image: Path to a background image to use (if no base_ppt).
    """
    
    # Extract dynamic theme color from LLM response
    dynamic_color_hex = presentation_data.get('theme_color', None)
    if dynamic_color_hex and dynamic_color_hex.startswith('#') and len(dynamic_color_hex) == 7:
        try:
             # Convert hex to RGB tuple
             hex_code = dynamic_color_hex.lstrip('#')
             theme_color = tuple(int(hex_code[i:i+2], 16) for i in (0, 2, 4))
        except:
             pass # Fallback to passed-in theme_color

    theme_rgb = RGBColor(*theme_color)
    
    if base_ppt and os.path.exists(base_ppt):
        try:
            prs = Presentation(base_ppt)
        except Exception as e:
            print(f"Error loading base PPT: {e}")
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

    if not presentation_data or 'slides' not in presentation_data:
        print("Error: Invalid data received")
        return prs

    slides_data = presentation_data['slides']
    today_date = datetime.now().strftime("%b %d, %Y")

    # Canvas-Based Layout Approach
    LAYOUT_BLANK = 6 

    for i, slide_data in enumerate(slides_data):
        layout_type = slide_data.get('layout', 'content')
        content = slide_data.get('content', [])
        image_query = slide_data.get('image_query')
        
        # Calculate content length
        full_text = " ".join(content)
        text_length = len(full_text)
        
        # Get dynamic font size
        start_font_size = get_smart_font_size(text_length, layout_type)
        
        # SMART LAYOUT SWITCHING
        # If content is very long, force full-width layout to avoid overflow
        # If using 'content' layout, we use start_font_size logic.
        
        total_chars = text_length
        if layout_type == 'image_right' and total_chars > 600:
            print(f"Notice: Content too long ({total_chars} chars) for Image Layout. Switching to Full Content.")
            layout_type = 'content'
            
        # 1. ALWAYS Add Blank Slide
        try:
            slide_layout = prs.slide_layouts[LAYOUT_BLANK]
        except:
            slide_layout = prs.slide_layouts[-1] 
            
        slide = prs.slides.add_slide(slide_layout)

        # 2. Apply Background
        if not base_ppt and background_image and os.path.exists(background_image):
             slide.shapes.add_picture(background_image, 0, 0, width=prs.slide_width, height=prs.slide_height)
        elif not base_ppt:
             bg_image_path = "static/bg.jpg"
             if os.path.exists(bg_image_path):
                 slide.shapes.add_picture(bg_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

        # 3. Draw Title (Always present, standard position)
        title_text = slide_data.get('title', 'Untitled')
        
        if i == 0:
            layout_type = 'title'

        if layout_type == 'title':
             # Centered Title (Middle of slide)
             # Use a large box for title
             title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9.0), Inches(3.0))
             title_tf = title_box.text_frame
             title_tf.word_wrap = True
             p = title_tf.paragraphs[0]
             p.text = title_text
             p.font.bold = True
             p.font.name = 'Arial'
             p.font.color.rgb = theme_rgb
             p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
             
             # Robust Font Sizing for Title
             if len(title_text) > 60:
                 p.font.size = Pt(32)
             elif len(title_text) > 30:
                 p.font.size = Pt(44)
             else:
                 p.font.size = Pt(54)
             
             # Subtitle if exists
             if content:
                 # Subtitle below title
                 sub_shape = slide.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(8.0), Inches(2.0))
                 # Use subtitle content (safely get first string)
                 if isinstance(content, list) and len(content) > 0:
                     subtitle_text = content[0]
                 elif isinstance(content, str):
                     subtitle_text = content
                 else:
                     subtitle_text = ""
                     
                 if subtitle_text:
                     fit_text(sub_shape.text_frame, [subtitle_text], font_size_pt=24, is_quote=True, font_family='Arial')
                 
        else:
             # Standard Top Title
             title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.4))
             title_tf = title_shape.text_frame
             title_tf.word_wrap = True
             p = title_tf.paragraphs[0]
             p.text = title_text
             
             # Dynamic Font Size
             if len(title_text) > 80:
                 p.font.size = Pt(28)
             elif len(title_text) > 50:
                 p.font.size = Pt(36)
             else:
                 p.font.size = Pt(40)
                 
             p.font.bold = True
             p.font.name = 'Arial'
             p.font.color.rgb = theme_rgb
             p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
             
             # Accent Line
             line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(9.0), Inches(0.05))
             line.fill.solid()
             line.fill.fore_color.rgb = theme_rgb
             line.line.fill.background()

             # 4. Handle Content Layouts
             try:
                 # Layout Logic
                 if layout_type == 'big_quote':
                     # Centered Text Box
                     quote_shape = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(8.0), Inches(4.0))
                     fit_text(quote_shape.text_frame, content, font_size_pt=start_font_size, is_quote=True, font_family='Arial')
                     
                 elif layout_type == 'image_right':
                     # Two Boxes: Text Left (5"), Image Right (3.5")
                     # Start Content at 2.0
                     
                     # Text Box (Left)
                     text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.0), Inches(5.0))
                     fit_text(text_box.text_frame, content, font_size_pt=start_font_size, font_family='Arial')
                     
                     # Image Box (Right)
                     if image_query:
                         print(f"Fetching image for: {image_query}")
                         img_path = fetch_image(image_query)
                         if img_path and os.path.exists(img_path):
                             try:
                                 # Add Picture at calculated position
                                 slide.shapes.add_picture(img_path, Inches(6.0), Inches(2.0), width=Inches(3.5), height=Inches(4.5))
                             except Exception as e:
                                 print(f"Error drawing image: {e}")
                                 
                 elif layout_type == 'two_columns': 
                     # Two Text Boxes
                     col1 = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4.2), Inches(5.0))
                     col2 = slide.shapes.add_textbox(Inches(5.3), Inches(2.0), Inches(4.2), Inches(5.0))
                     
                     mid = len(content) // 2
                     fit_text(col1.text_frame, content[:mid], font_size_pt=start_font_size, font_family='Arial')
                     fit_text(col2.text_frame, content[mid:], font_size_pt=start_font_size, font_family='Arial')

                 elif layout_type == 'chart':
                     # Two Boxes: Text Left (3.5"), Chart Right (6")
                     
                     # Text Box (Left)
                     text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(3.5), Inches(5.0))
                     fit_text(text_box.text_frame, content, font_size_pt=start_font_size, font_family='Arial')
                     
                     # Chart Box (Right)
                     try:
                         chart_data_dict = slide_data.get('chart_data')
                         if chart_data_dict:
                             # Build chart data
                             chart_data = CategoryChartData()
                             chart_data.categories = chart_data_dict.get('categories', [])
                             for series in chart_data_dict.get('series', []):
                                 chart_data.add_series(series.get('name', 'Series'), series.get('values', []))
                             
                             # Determine Chart Type
                             chart_type_str = chart_data_dict.get('type', 'bar')
                             if chart_type_str == 'pie': chart_type = XL_CHART_TYPE.PIE
                             elif chart_type_str == 'line': chart_type = XL_CHART_TYPE.LINE
                             else: chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
                             
                             # Add chart
                             x, y, cx, cy = Inches(4.2), Inches(2.0), Inches(5.3), Inches(4.5)
                             chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart
                             
                             # Professional Consulting Chart Styling
                             chart.has_legend = True
                             chart.legend.include_in_layout = False
                             chart.legend.position = XL_LEGEND_POSITION.TOP
                             
                             # Add Chart Title if provided
                             c_title = chart_data_dict.get('chart_title')
                             if c_title:
                                 chart.has_title = True
                                 chart.chart_title.text_frame.text = c_title
                                 chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
                                 chart.chart_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(60, 60, 60)
                             
                             # Clean up axes (remove clutter)
                             if chart_type in (XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.LINE):
                                 val_axis = chart.value_axis
                                 val_axis.has_major_gridlines = True # Keep horizontal bounds but make them light
                                 val_axis.major_tick_mark = XL_TICK_MARK.NONE
                                 val_axis.minor_tick_mark = XL_TICK_MARK.NONE
                                 cat_axis = chart.category_axis
                                 cat_axis.major_tick_mark = XL_TICK_MARK.NONE
                             
                             # Enable data labels for precision
                             try:
                                 for plot in chart.plots:
                                     plot.has_data_labels = True
                                     data_labels = plot.data_labels
                                     data_labels.font.size = Pt(10)
                                     data_labels.font.color.rgb = RGBColor(80, 80, 80)
                             except Exception as dle:
                                 pass
                             
                             # Color the series with variations of the theme color
                             try:
                                  # Extract base RGB values
                                  r, g, b = theme_rgb
                                  
                                  for idx, series in enumerate(chart.series):
                                      fill = series.format.fill
                                      fill.solid()
                                      # Lighten the color for subsequent series to create a nice palette
                                      factor = 1.0 + (idx * 0.4) 
                                      new_r = min(255, int(r * factor)) if idx > 0 else r
                                      new_g = min(255, int(g * factor)) if idx > 0 else g
                                      new_b = min(255, int(b * factor)) if idx > 0 else b
                                      
                                      # If it gets completely white, fallback to a grey
                                      if new_r == 255 and new_g == 255 and new_b == 255:
                                          fill.fore_color.rgb = RGBColor(200, 200, 200)
                                      else:
                                          fill.fore_color.rgb = RGBColor(new_r, new_g, new_b)
                             except: pass
                         else:
                             print("Warning: Layout is chart but chart_data is missing.")
                     except Exception as ce:
                         print(f"Error drawing chart: {ce}")
     
                 else: # 'content' or default
                     # Single Full Width Box
                     text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9.0), Inches(5.0))
                     fit_text(text_box.text_frame, content, font_size_pt=start_font_size, font_family='Arial')
             except Exception as e:
                 print(f"Error rendering content: {e}")
                 # Fallback to simple text
                 try:
                     tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9.0), Inches(5.0))
                     tb.text_frame.text = "Error rendering content. Please check logs."
                 except: pass

        # Footer (Page number)

        try:
             slide_num_box = slide.shapes.add_textbox(Inches(9.0), Inches(7.0), Inches(0.8), Inches(0.4))
             p = slide_num_box.text_frame.paragraphs[0]
             p.text = f"{i+1}"
             p.font.size = Pt(12)
             p.font.color.rgb = RGBColor(128, 128, 128)
             p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
             
             date_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(2.0), Inches(0.4))
             p = date_box.text_frame.paragraphs[0]
             p.text = today_date
             p.font.size = Pt(12)
             p.font.color.rgb = RGBColor(128, 128, 128)
        except Exception:
             pass

        # Speaker Notes
        if 'speaker_notes' in slide_data:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data['speaker_notes']

    return prs
