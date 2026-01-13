import os
from io import BytesIO
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE

def pptx_to_markdown(pptx_bytes, output_dir="extracted_images"):
    prs = Presentation(BytesIO(pptx_bytes))
    os.makedirs(output_dir, exist_ok=True)
    markdown_lines = []
    for i, slide in enumerate(prs.slides, start=1):
        fill = slide.background.fill
        if fill.type == MSO_FILL.PICTURE:
            blip = fill._xFill.blipFill.blip
            if blip is not None:
                rId = blip.rEmbed
                image_part = slide.part.related_part(rId)
                bg_filename = f"{output_dir}/slide{i}_bg.jpeg"
                with open(bg_filename, "wb") as f:
                    f.write(image_part.blob)
                markdown_lines.append(f"![Background]({os.path.basename(bg_filename)})\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                if not markdown_lines[-1].startswith("##") and not markdown_lines[-1].startswith("# "):
                    markdown_lines.append(f"## {text}\n")
                else:
                    markdown_lines.append(text + "\n")

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                img_filename = f"{output_dir}/slide{i}_{shape.shape_id}.{ext}"
                with open(img_filename, "wb") as f:
                    f.write(image.blob)
                markdown_lines.append(f"![Image]({os.path.basename(img_filename)})\n")

        markdown_lines.append("\n---\n")
    return "\n".join(markdown_lines)
